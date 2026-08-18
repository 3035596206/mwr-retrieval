# -*- coding: utf-8 -*-
"""期末汇报 PPT 生成脚本 —— 地基微波辐射计大气温湿廓线反演项目。

面向课题组例会的精简版 (~13 页)，覆盖 BRNN 统计反演 / OEM 物理反演 /
成都真实数据验证 / 整体进度与计划 四条主线，嵌入项目现有结果图。

用法:
    python scripts/make_final_report_ppt.py
产物:
    reports/汇报PPT/期末汇报_MWR温湿廓线反演_2026-07-29.pptx
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ----------------------------------------------------------------------
# 路径与常量
# ----------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parent.parent
RESULTS = PROJECT_ROOT / "results"
OUT_DIR = PROJECT_ROOT / "reports" / "汇报PPT"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "期末汇报_MWR温湿廓线反演_2026-07-29.pptx"

# 16:9 画布 (PowerPoint 默认宽屏)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色：深蓝学术风
C_PRIMARY = RGBColor(0x0B, 0x3D, 0x91)      # 主色 深蓝
C_PRIMARY_DARK = RGBColor(0x07, 0x28, 0x5E) # 更深
C_ACCENT = RGBColor(0xE8, 0x7A, 0x1E)       # 强调 橙
C_LIGHT = RGBColor(0xEE, 0xF2, 0xF9)        # 浅底
C_TEXT = RGBColor(0x1F, 0x2A, 0x37)         # 正文
C_MUTED = RGBColor(0x5A, 0x6B, 0x7B)        # 次要文字
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE = RGBColor(0xD0, 0xD9, 0xE5)

FONT_CN = "微软雅黑"
FONT_EN = "Calibri"


# ----------------------------------------------------------------------
# 底层辅助
# ----------------------------------------------------------------------
def _set_run(run, text, size, bold=False, color=C_TEXT, font=FONT_CN,
             italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    # 中文字体名需写入 east-asian 属性才会生效
    rPr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_CN)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=C_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font=FONT_CN, line_spacing=1.15, italic=False):
    """单段文本框。text 可含 \n。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        _set_run(run, line, size, bold=bold, color=color, font=font,
                 italic=italic)
    return tb


def add_bullets(slide, left, top, width, height, items, size=14,
                color=C_TEXT, line_spacing=1.25, bullet_color=C_ACCENT):
    """多级要点列表。item = (text, level) 或 text。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        indent = "    " * level
        marker = "● " if level == 0 else "– "
        r1 = p.add_run()
        _set_run(r1, indent + marker, size, bold=False, color=bullet_color)
        r2 = p.add_run()
        _set_run(r2, text, size, bold=False, color=color)
    return tb


def add_rect(slide, left, top, width, height, fill, line=None,
             line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def add_header(slide, title, page_no, total, subtitle=None):
    """统一的页眉：左侧色条 + 标题 + 页码。"""
    # 顶部主色条
    add_rect(slide, Emu(0), Emu(0), Inches(0.22), Inches(7.5), C_PRIMARY)
    # 标题
    add_textbox(slide, Inches(0.55), Inches(0.30), Inches(11.5), Inches(0.55),
                title, size=24, bold=True, color=C_PRIMARY_DARK,
                anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.86), Inches(11.5),
                    Inches(0.3), subtitle, size=12, color=C_MUTED)
    # 标题下分隔线
    add_rect(slide, Inches(0.58), Inches(1.18), Inches(12.2), Emu(20000),
             C_LINE)
    # 页码
    add_textbox(slide, Inches(11.9), Inches(7.05), Inches(1.3), Inches(0.3),
                f"{page_no} / {total}", size=10, color=C_MUTED,
                align=PP_ALIGN.RIGHT)
    # 底部页脚
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(10), Inches(0.3),
                "地基微波辐射计大气温湿廓线反演  ·  期末汇报",
                size=10, color=C_MUTED)


def add_table(slide, left, top, width, height, data, col_widths=None,
              header_fill=C_PRIMARY, header_color=C_WHITE, font_size=12,
              header_size=12, zebra=True):
    """data: 二维 list，第一行为表头。返回表格 shape。"""
    rows, cols = len(data), len(data[0])
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gtable.table
    # 关闭默认带状样式
    tbl = table._tbl
    from pptx.oxml.ns import qn
    for child in tbl.findall(qn("a:tblPr")):
        child.set("firstRow", "0")
        child.set("bandRow", "0")
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(22860)
            cell.margin_bottom = Emu(22860)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            text = str(data[r][c])
            if r == 0:
                _set_run(run, text, header_size, bold=True, color=header_color)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                bold = (c == 0)
                color = C_TEXT
                # 强调数值列关键字
                if any(k in text for k in ("★", "当前最佳", "推荐")):
                    bold = True
                    color = C_ACCENT
                _set_run(run, text, font_size, bold=bold, color=color)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (C_LIGHT if zebra and r % 2 == 0
                                            else C_WHITE)
    return gtable


def add_image_fit(slide, path, left, top, max_w, max_h, align="center",
                  valign="middle"):
    """按比例缩放图片放入 max_w x max_h 区域，返回最终绘图区信息。"""
    path = Path(path)
    if not path.exists():
        add_textbox(slide, left, top, max_w, max_h,
                    f"[图片缺失] {path.name}", size=12, color=C_ACCENT,
                    anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * scale))
    h = Emu(int(ih * scale))
    if align == "center":
        x = left + (max_w - w) // 2
    elif align == "right":
        x = left + (max_w - w)
    else:
        x = left
    if valign == "middle":
        y = top + (max_h - h) // 2
    elif valign == "bottom":
        y = top + (max_h - h)
    else:
        y = top
    slide.shapes.add_picture(str(path), x, y, w, h)


def add_caption(slide, left, top, width, text):
    add_textbox(slide, left, top, width, Inches(0.3), text, size=10,
                color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


def card(slide, left, top, width, height, title, lines, title_color=C_WHITE,
         title_fill=C_PRIMARY, body_fill=C_LIGHT, t_size=13, b_size=11):
    """一个小卡片：标题条 + 正文要点。"""
    bar_h = Inches(0.38)
    add_rect(slide, left, top, width, bar_h, title_fill)
    add_rect(slide, left, top + bar_h, width, height - bar_h, body_fill)
    add_textbox(slide, left + Inches(0.12), top, width - Inches(0.2), bar_h,
                title, size=t_size, bold=True, color=title_color,
                anchor=MSO_ANCHOR.MIDDLE)
    if lines:
        add_bullets(slide, left + Inches(0.14), top + bar_h + Inches(0.08),
                    width - Inches(0.24), height - bar_h - Inches(0.16),
                    lines, size=b_size, line_spacing=1.18)


# ======================================================================
# 页面构建
# ======================================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    TOTAL = 13

    # ---------------- 1. 封面 ----------------
    s = prs.slides.add_slide(blank)
    # 背景大色块
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, C_PRIMARY_DARK)
    add_rect(s, Emu(0), Inches(4.55), SLIDE_W, Inches(0.06), C_ACCENT)
    add_rect(s, Emu(0), Inches(0), Inches(0.35), SLIDE_H, C_ACCENT)
    add_textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
                "期末汇报  ·  课题组例会", size=16, color=RGBColor(0xBD,0xD4,0xF0))
    add_textbox(s, Inches(0.9), Inches(2.05), Inches(11.6), Inches(1.7),
                "地基微波辐射计\n大气温湿廓线反演", size=44, bold=True,
                color=C_WHITE, line_spacing=1.1)
    add_textbox(s, Inches(0.92), Inches(4.75), Inches(11.5), Inches(0.5),
                "BRNN 统计反演  ·  OEM 物理反演  ·  成都真实数据验证",
                size=18, color=RGBColor(0xE6,0xED,0xF7))
    # 关键指标条
    metrics = [
        ("BRNN v4", "T 1.26 K / RH 7.76%"),
        ("OEM MonoRTM n=100", "T 2.02 K · DOFS 2.21"),
        ("成都 21ch", "真实亮温链路打通"),
    ]
    for i, (k, v) in enumerate(metrics):
        x = Inches(0.9 + i * 3.9)
        add_textbox(s, x, Inches(5.5), Inches(3.8), Inches(0.4), k,
                    size=13, bold=True, color=C_ACCENT)
        add_textbox(s, x, Inches(5.9), Inches(3.8), Inches(0.4), v,
                    size=13, color=C_WHITE)
    add_textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
                "汇报日期：2026-07-29", size=13,
                color=RGBColor(0x9F,0xB6,0xD6))

    # ---------------- 2. 研究背景与目标 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "研究背景与目标", 2, TOTAL)
    card(s, Inches(0.58), Inches(1.45), Inches(6.05), Inches(2.55),
         "▍ 研究背景", [
             "大气温湿廓线是数值预报、强对流与航空气象的关键输入。",
             "探空精度高但时空稀疏（每日 2 次）；微波辐射计(MWR)可实现高频连续探测。",
             "反演难点：欠定问题，垂直信息量有限，云天/降雨干扰大。",
         ])
    card(s, Inches(6.78), Inches(1.45), Inches(6.0), Inches(2.55),
         "▍ 项目目标", [
             "建立地基多通道 MWR 温度 T(z) 与相对湿度 RH(z) 反演系统。",
             "复现并超越论文基准（T<1.5K, RH<13%）。",
             "构建 BRNN 统计反演 + OEM 物理反演混合框架。",
             "在成都本地真实数据上完成验证。",
         ])
    # 数据与通道
    add_textbox(s, Inches(0.58), Inches(4.25), Inches(12.2), Inches(0.35),
                "▍ 数据与仪器", size=15, bold=True, color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(4.65), Inches(12.2), Inches(2.0), [
        ["数据 / 仪器", "通道配置", "覆盖", "用途"],
        ["MP-3000A（北京南郊）", "22 通道 K/V 波段", "2013–2019", "BRNN v4 训练（历史最佳）"],
        ["RPG HATPRO 配置", "14 通道（K7+V7）", "2013-01 ERA5", "OEM/MonoRTM 物理反演主线"],
        ["成都实测辐射计", "21 通道（K7+V7+W1+G5+高频1）", "2026-05", "真实数据验证"],
        ["温江探空", "483 个文件", "2026-05", "独立外部验证真值"],
    ], col_widths=[Inches(3.0), Inches(4.3), Inches(1.9), Inches(3.0)])

    # ---------------- 3. 技术路线总览 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "技术路线总览：统计 + 物理混合反演", 3, TOTAL)
    # 流程框
    boxes = [
        (0.58, "观测亮温\nObs_BT", C_PRIMARY),
        (2.78, "QC / 偏差订正", C_PRIMARY),
    ]
    y_flow = Inches(1.7)
    bw, bh = Inches(2.0), Inches(1.2)
    # 上分支 BRNN
    add_rect(s, Inches(0.58), y_flow, bw, bh, C_PRIMARY)
    add_textbox(s, Inches(0.58), y_flow, bw, bh, "观测亮温\nObs_BT",
                size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(2.78), y_flow, bw, bh, C_PRIMARY)
    add_textbox(s, Inches(2.78), y_flow, bw, bh, "QC / 偏差订正\n+ ERA5 标签",
                size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # 分叉到 BRNN / OEM
    add_rect(s, Inches(5.0), Inches(1.35), Inches(2.2), Inches(0.95), C_ACCENT)
    add_textbox(s, Inches(5.0), Inches(1.35), Inches(2.2), Inches(0.95),
                "BRNN 统计反演\n~0.5 ms/廓线",
                size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(5.0), Inches(2.55), Inches(2.2), Inches(0.95),
             C_ACCENT)
    add_textbox(s, Inches(5.0), Inches(2.55), Inches(2.2), Inches(0.95),
                "OEM 物理反演\n(1D-Var)",
                size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # 右侧产物
    add_rect(s, Inches(7.4), Inches(1.35), Inches(5.35), Inches(0.95),
             C_PRIMARY_DARK)
    add_textbox(s, Inches(7.55), Inches(1.35), Inches(5.1), Inches(0.95),
                "T(z), RH(z) 廓线  ·  6 子模型  ·  当前最佳 v4",
                size=12, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(7.4), Inches(2.55), Inches(5.35), Inches(0.95),
             C_PRIMARY_DARK)
    add_textbox(s, Inches(7.55), Inches(2.55), Inches(5.1), Inches(0.95),
                "T/RH + 后验误差 + AK + DOFS  ·  ARTS 主前向模型",
                size=12, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 三层能力
    add_textbox(s, Inches(0.58), Inches(4.05), Inches(12.2), Inches(0.35),
                "▍ 已形成的三层能力", size=15, bold=True, color=C_PRIMARY_DARK)
    card(s, Inches(0.58), Inches(4.45), Inches(3.95), Inches(2.4),
         "第一层 · 统计反演", [
             "BRNN v4 当前最佳",
             "T = 1.26 K，RH = 7.76%",
             "推理 ~0.5 ms/廓线",
         ], title_fill=C_PRIMARY)
    card(s, Inches(4.69), Inches(4.45), Inches(3.95), Inches(2.4),
         "第二层 · 物理反演", [
             "OEM + ARTS/MonoRTM",
             "n=100 self-consistent",
             "DOFS ≈ 2.21，收敛 99%",
         ], title_fill=C_ACCENT)
    card(s, Inches(8.8), Inches(4.45), Inches(3.95), Inches(2.4),
         "第三层 · 混合路线", [
             "NN 作先验/S_a",
             "NN surrogate 加速 H(x)",
             "不确定度校准 + 云天 OEM",
         ], title_fill=C_PRIMARY_DARK)

    # ---------------- 4. BRNN 统计反演方法 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "BRNN 统计反演：方法与 v4 关键改进", 4, TOTAL)
    add_textbox(s, Inches(0.58), Inches(1.4), Inches(12.2), Inches(0.35),
                "▍ 网络结构：6 个独立 BRNN 子模型（T/RH × 0-2km / 2-8km / 8-10km）",
                size=14, bold=True, color=C_PRIMARY_DARK)
    add_bullets(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(1.3), [
        "输入：多通道亮温 + 地面 T2m/RH2m/Ps/IR/CLWC/CIC 等辅助特征；输出 0-10km 共 93 层 T/RH。",
        "分高度段建模：各高度区间物理特性差异大，独立网络更稳定；边界层(2km/8km)由相邻模型各预测一次。",
        "Hidden=256, Dropout=0.3, LR=1e-3, Batch=128, Early-stopping patience=20。",
    ], size=12.5)
    add_textbox(s, Inches(0.58), Inches(3.15), Inches(12.2), Inches(0.35),
                "▍ v4 相对前序版本的关键改进（成就当前最佳）", size=14,
                bold=True, color=C_PRIMARY_DARK)
    items = [
        ("CLWC 廓线筛除 + K 波段逐通道 2.5σ 筛除", "剔除云天/异常亮温样本"),
        ("Winsorize 回归 BT 订正（非 OLS）", "K 波段 OMB 厚尾 → OMB std −49%"),
        ("IR_Temperature 入模", "引入红外测温作为辅助特征"),
        ("按 Profile_Index 分组划分（非时间划分）", "同一廓线重复观测~5.6 次，杜绝数据泄露"),
    ]
    y0 = Inches(3.6)
    for i, (t, d) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.58 + col * 6.1)
        yy = y0 + Inches(row * 1.0)
        add_rect(s, x, yy, Inches(0.12), Inches(0.85), C_ACCENT)
        add_textbox(s, x + Inches(0.25), yy, Inches(5.7), Inches(0.45), t,
                    size=12.5, bold=True, color=C_TEXT)
        add_textbox(s, x + Inches(0.25), yy + Inches(0.45), Inches(5.7),
                    Inches(0.4), "→ " + d, size=11, color=C_MUTED)
    add_textbox(s, Inches(0.58), Inches(6.5), Inches(12.2), Inches(0.35),
                "核心结论：v3 证实 Sim→Obs domain gap 退化 3.7×，故 v4 改为 Obs_BT 直接训练，"
                "成为当前最优统计反演模型。",
                size=12, color=C_PRIMARY_DARK, italic=True)

    # ---------------- 5. BRNN 结果 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "BRNN v4 反演结果（MP-3000A）", 5, TOTAL)
    # 左：指标表
    add_textbox(s, Inches(0.58), Inches(1.4), Inches(4.3), Inches(0.35),
                "▍ 核心指标", size=14, bold=True, color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(1.8), Inches(4.35), Inches(1.7), [
        ["指标", "v4", "论文基准"],
        ["T RMSE", "1.26 K", "<1.5 K"],
        ["RH RMSE", "7.76%", "<13%"],
        ["推理速度", "~0.5 ms/廓线", "—"],
    ], col_widths=[Inches(1.5), Inches(1.45), Inches(1.4)], font_size=12.5)
    add_textbox(s, Inches(0.58), Inches(3.75), Inches(4.3), Inches(0.35),
                "▍ 分层精度", size=14, bold=True, color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(4.15), Inches(4.35), Inches(2.2), [
        ["高度区间", "T RMSE", "RH RMSE"],
        ["0–0.5 km", "1.09 K", "5.78%"],
        ["2–8 km", "1.30 K", "9.50%"],
        ["5–10 km", "1.50 K", "10.31%"],
    ], col_widths=[Inches(1.7), Inches(1.35), Inches(1.3)], font_size=12)
    add_textbox(s, Inches(0.58), Inches(6.45), Inches(4.35), Inches(0.6),
                "近地层精度最优；高空信息量下降，与 MWR 垂直分辨能力一致。",
                size=10.5, color=C_MUTED, italic=True)
    # 右：散点图 + 误差廓线
    add_image_fit(s, RESULTS / "T_scatter.png", Inches(5.2), Inches(1.4),
                  Inches(3.8), Inches(2.7))
    add_caption(s, Inches(5.2), Inches(4.05), Inches(3.8),
                "图1  温度散点（v4，密集区接近 1:1）")
    add_image_fit(s, RESULTS / "RH_error_profile.png", Inches(9.15),
                  Inches(1.4), Inches(3.8), Inches(2.7))
    add_caption(s, Inches(9.15), Inches(4.05), Inches(3.8),
                "图2  相对湿度误差廓线")
    add_image_fit(s, RESULTS / "T_error_profile.png", Inches(5.2),
                  Inches(4.45), Inches(3.8), Inches(2.5))
    add_caption(s, Inches(5.2), Inches(6.95), Inches(3.8),
                "图3  温度误差廓线")

    # ---------------- 6. BRNN 迭代历程 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "BRNN 迭代历程与关键决策", 6, TOTAL)
    add_textbox(s, Inches(0.58), Inches(1.4), Inches(12.2), Inches(0.35),
                "▍ 六轮迭代：从数据泄露基线到当前最佳 v4", size=14, bold=True,
                color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(1.8), Inches(12.2), Inches(2.9), [
        ["版本", "方案", "T RMSE", "RH RMSE", "结论"],
        ["v1", "原始 Obs_BT，时间划分", "3.03 K", "16.55%", "基线，存在数据泄露与质量问题"],
        ["v2", "Obs_BT 订正 + 廓线分组", "1.45 K", "9.03%", "首个强结果"],
        ["v3", "Sim_BT 训练，Obs_BT 测试", "2.65 K", "12.00%", "暴露 Sim→Obs domain gap"],
        ["v4 ★", "Obs_BT 订正 + K波段过滤 + IR入模", "1.26 K", "7.76%", "当前最佳"],
        ["v6", "Sim→Obs 两阶段训练", "1.92 K", "10.34%", "Sim 路线有进展但未超 v4"],
    ], col_widths=[Inches(1.1), Inches(4.6), Inches(1.5), Inches(1.5),
                   Inches(3.5)], font_size=12)
    add_textbox(s, Inches(0.58), Inches(5.0), Inches(12.2), Inches(0.35),
                "▍ 关键决策记录", size=14, bold=True, color=C_PRIMARY_DARK)
    decisions = [
        ("Obs_BT 直接训练（v4）而非 Sim_BT", "v3 证明 Sim→Obs 退化 3.7×"),
        ("6 个独立 BRNN 而非单一多输出", "各高度段物理特性差异大"),
        ("廓线分组划分而非时间划分", "同廓线重复观测~5.6次，必须防泄露"),
        ("Winsorize 回归订正而非 OLS", "K 波段 OMB 厚尾分布"),
    ]
    for i, (t, d) in enumerate(decisions):
        col = i % 2
        row = i // 2
        x = Inches(0.58 + col * 6.1)
        yy = Inches(5.45 + row * 0.8)
        add_textbox(s, x, yy, Inches(5.9), Inches(0.35),
                    "▸ " + t, size=12, bold=True, color=C_TEXT)
        add_textbox(s, x + Inches(0.3), yy + Inches(0.33), Inches(5.6),
                    Inches(0.35), d, size=10.5, color=C_MUTED, italic=True)

    # ---------------- 7. OEM 物理反演框架 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "OEM 物理反演框架", 7, TOTAL)
    add_textbox(s, Inches(0.58), Inches(1.4), Inches(12.2), Inches(0.35),
                "▍ 最优估计 (OEM / 1D-Var) 目标函数", size=14, bold=True,
                color=C_PRIMARY_DARK)
    # 公式框
    add_rect(s, Inches(0.58), Inches(1.85), Inches(12.2), Inches(0.75),
             C_LIGHT)
    add_textbox(s, Inches(0.58), Inches(1.85), Inches(12.2), Inches(0.75),
                "J(x) = (x−xₐ)ᵀ Sₐ⁻¹ (x−xₐ)  +  (y−H(x))ᵀ Sₑ⁻¹ (y−H(x))",
                size=17, bold=True, color=C_PRIMARY_DARK, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    # 符号表
    add_table(s, Inches(0.58), Inches(2.8), Inches(7.4), Inches(3.3), [
        ["符号", "含义", "当前实现"],
        ["x", "状态向量", "14d T7+RH7（→ 21d 含 LWC）"],
        ["xₐ", "背景场", "ERA5 扰动（→ BRNN first guess）"],
        ["Sₐ", "背景误差协方差", "指数相关 / v4-derived"],
        ["y", "观测亮温", "synthetic / MonoRTM / Obs_BT"],
        ["H(x)", "前向模型", "ARTS 主线；MonoRTM 历史"],
        ["Sₑ", "观测误差协方差", "K-band 1.5K, V-band 0.5K"],
    ], col_widths=[Inches(1.1), Inches(2.4), Inches(3.9)], font_size=11.5)
    # 右：能力清单
    card(s, Inches(8.18), Inches(2.8), Inches(4.6), Inches(3.7),
         "▍ 已实现的 OEM 能力", [
             "LM / Gauss-Newton 求解器",
             "有限差分 Jacobian",
             "状态向量打包/解包",
             "Averaging kernel / DOFS",
             "后验协方差诊断",
             "self-consistent 闭环验证",
         ], title_fill=C_ACCENT, b_size=11.5)
    add_textbox(s, Inches(0.58), Inches(6.35), Inches(12.2), Inches(0.6),
                "状态向量粗分层：0-0.5 / 0.5-1 / 1-2 / 2-3 / 3-5 / 5-8 / 8-10 km，"
                "T7+RH7 共 14 维；后续扩展 10d EOF 与 21d 云天状态。",
                size=11, color=C_MUTED, italic=True)

    # ---------------- 8. OEM 实验结果 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "OEM 实验结果：MonoRTM n=100 基线", 8, TOTAL)
    add_textbox(s, Inches(0.58), Inches(1.4), Inches(7.2), Inches(0.35),
                "▍ 2013-01 三类 POC + n=100 基线", size=14, bold=True,
                color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(1.8), Inches(7.2), Inches(2.4), [
        ["实验（前向/样本）", "T RMSE", "RH RMSE", "BT RMS", "DOFS"],
        ["Self-consistent simple n=20", "2.64→1.95K", "6.43→6.06%", "1.59→0.54K", "2.10"],
        ["Forward-mismatch simple n=20", "→27.0K", "7.59→7.22%", "→20.2K", "2.21"],
        ["MonoRTM self-consistent n=100", "2.64→2.02K", "6.49→6.20%", "5.03→0.61K", "2.21"],
    ], col_widths=[Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.1),
                   Inches(0.7)], font_size=10.5, header_size=10.5)
    add_textbox(s, Inches(0.58), Inches(4.35), Inches(7.2), Inches(0.35),
                "▍ 核心结论", size=13, bold=True, color=C_PRIMARY_DARK)
    add_bullets(s, Inches(0.7), Inches(4.75), Inches(7.0), Inches(2.0), [
        "self-consistent 验证算法链路可用，收敛率 99%、平均迭代 7.79。",
        "forward-mismatch 证明 H(x) 与观测不一致会导致 T 严重退化 → 必须先闭环验证。",
        "MonoRTM 是当前最可信物理 POC；BT residual 降幅最大（5.03→0.61K）。",
        "DOFS≈2.21/14，符合地基 MWR 垂直信息量有限的预期。",
    ], size=11.5)
    # 右图：改进汇总
    add_image_fit(s, RESULTS / "oem_201301_self_consistent_monortm_n100"
                  /"rmse_profiles.png", Inches(8.0), Inches(1.5),
                  Inches(4.85), Inches(2.7))
    add_caption(s, Inches(8.0), Inches(4.2), Inches(4.85),
                "图4  MonoRTM n=100：Prior→Posterior RMSE 廓线")
    add_image_fit(s, RESULTS / "oem_201301_self_consistent_monortm_n100"
                  /"bt_dofs.png", Inches(8.0), Inches(4.55), Inches(4.85),
                  Inches(2.3))
    add_caption(s, Inches(8.0), Inches(6.85), Inches(4.85),
                "图5  BT 残差收敛 与 DOFS 分布")

    # ---------------- 9. ARTS / 诊断 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "ARTS 前向模型与垂直信息量诊断", 9, TOTAL)
    card(s, Inches(0.58), Inches(1.4), Inches(5.95), Inches(1.85),
         "▍ 前向模型主线切换", [
             "ARTS 设为默认 OEM 后端（对齐研究组工作流）。",
             "支持成都 21 通道；通过 WSL pyarts runner 接入本地 agenda。",
             "MonoRTM 保留为历史基线，用于交叉验证。",
         ], b_size=11.5)
    card(s, Inches(6.68), Inches(1.4), Inches(6.1), Inches(1.85),
         "▍ 关键诊断结论", [
             "OEM 改善集中在 0-2km 近地层；5-10km T 反演几乎完全依赖先验。",
             "DOFS 仅 2.2/14 → 高空信息量极低，是核心瓶颈。",
             "需通过多仰角观测 / 高频通道 / 更优 Sₐ 提升信息量。",
         ], title_fill=C_ACCENT, b_size=11.5)
    # 分层精度表
    add_textbox(s, Inches(0.58), Inches(3.45), Inches(5.95), Inches(0.35),
                "▍ BRNN v4 vs OEM 分层精度", size=13, bold=True,
                color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(3.85), Inches(5.95), Inches(1.8), [
        ["高度区间", "BRNN v4 Δ", "OEM ΔT / ΔRH"],
        ["0–0.5 km", "T1.09K RH5.78%", "+0.42K / +0.67%"],
        ["2–8 km", "T1.30K RH9.50%", "+0.07K / +0.28%"],
        ["5–10 km", "T1.50K RH10.31%", "−0.01K / +0.10%"],
    ], col_widths=[Inches(1.5), Inches(2.25), Inches(2.2)], font_size=10.5,
              header_size=10.5)
    # S_a / EOF
    add_textbox(s, Inches(0.58), Inches(5.85), Inches(5.95), Inches(0.35),
                "▍ Sₐ 与状态降维探索", size=13, bold=True, color=C_PRIMARY_DARK)
    add_bullets(s, Inches(0.7), Inches(6.25), Inches(5.8), Inches(1.0), [
        "v4-derived Sₐ (14×14) 已生成；T 改善 +22%。",
        "EOF/PCA：T/RH 各 5 EOF，基础就绪，10d 版本待优化。",
    ], size=11)
    # 右图：AK + S_a 对比
    add_image_fit(s, RESULTS / "oem_201301_self_consistent_monortm_n100"
                  /"averaging_kernel.png", Inches(6.7), Inches(3.4),
                  Inches(3.0), Inches(3.4))
    add_caption(s, Inches(6.7), Inches(6.85), Inches(3.0),
                "图6  Averaging Kernel")
    add_image_fit(s, RESULTS / "oem_covariance" / "sa_comparison.png",
                  Inches(9.85), Inches(3.4), Inches(2.95), Inches(3.4))
    add_caption(s, Inches(9.85), Inches(6.85), Inches(2.95),
                "图7  Sₐ 对比（指数 vs v4）")

    # ---------------- 10. 成都验证（一）ERA5 配对 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "成都真实数据验证（一）：ERA5 精确配对 Hybrid",
               10, TOTAL)
    add_textbox(s, Inches(0.58), Inches(1.35), Inches(12.2), Inches(0.35),
                "▍ 链路打通：21 通道实测亮温 + ERA5 精确小时标签",
                size=14, bold=True, color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(1.75), Inches(6.1), Inches(1.9), [
        ["数据集划分", "样本/天数"],
        ["实测亮温总记录", "163 条"],
        ["与 ERA5 同 UTC 配对", "139 条 / 16 天"],
        ["Train / Val / Test", "69 / 22 / 48 条"],
    ], col_widths=[Inches(3.4), Inches(2.7)], font_size=11.5)
    add_textbox(s, Inches(0.58), Inches(3.8), Inches(6.1), Inches(0.35),
                "▍ 当前最佳 Hybrid（ERA5 测试集）", size=13, bold=True,
                color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(4.2), Inches(6.1), Inches(2.0), [
        ["变量 / 方法", "RMSE", "相对气候态"],
        ["T · Ridge+5EOF (21ch)", "1.479 K", "改善 35.2%"],
        ["RH · BRNN 4-seed 集成", "20.966%", "改善 13.1%"],
    ], col_widths=[Inches(3.0), Inches(1.5), Inches(1.6)], font_size=11.5)
    add_textbox(s, Inches(0.58), Inches(6.3), Inches(6.1), Inches(0.7),
                "T=Ridge/EOF 更适合小样本；RH=BRNN 集成更适合非线性映射。"
                "首次实现 T/RH 同时优于气候态基线。",
                size=10.5, color=C_MUTED, italic=True)
    # 右图
    add_image_fit(s, RESULTS / "chengdu_era5_figures" / "03_prediction_scatter.png",
                  Inches(6.95), Inches(1.5), Inches(5.9), Inches(2.7))
    add_caption(s, Inches(6.95), Inches(4.2), Inches(5.9),
                "图8  成都 21ch Hybrid 预测散点（T/RH）")
    add_image_fit(s, RESULTS / "chengdu_era5_figures" / "02_height_resolved_errors.png",
                  Inches(6.95), Inches(4.55), Inches(5.9), Inches(2.4))
    add_caption(s, Inches(6.95), Inches(6.95), Inches(5.9),
                "图9  分高度误差（各通道组/模型对比）")

    # ---------------- 11. 成都验证（二）物理48层 + 探空 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "成都真实数据验证（二）：物理 48 层 + 温江探空独立验证",
               11, TOTAL)
    add_textbox(s, Inches(0.58), Inches(1.35), Inches(7.2), Inches(0.35),
                "▍ 改进：48 点 → 真正的 48 物理层 + 层平均保守映射",
                size=13.5, bold=True, color=C_PRIMARY_DARK)
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(7.0), Inches(1.35), [
        "层定义：0-500m×1 + 500-2000m/100m×15 + 2000-10000m/250m×32 = 48 层。",
        "地下气压层过滤（站点海拔 548m，均删 ~2.7 层）；T/log(q)/log(P) 层平均后重算 RH。",
        "选择性偏差订正：T 用低自由度探空订正（RMSE 1.067→0.656K）；RH 不订正。",
    ], size=11)
    add_textbox(s, Inches(0.58), Inches(3.15), Inches(7.2), Inches(0.35),
                "▍ 独立温江探空验证（6 份探空，1-2h 时间差）",
                size=13.5, bold=True, color=C_PRIMARY_DARK)
    add_table(s, Inches(0.58), Inches(3.55), Inches(7.2), Inches(2.0), [
        ["方案", "T RMSE", "T Bias", "RH RMSE", "RH Bias"],
        ["原 93 层（层平均）", "1.730K", "+0.998K", "23.685%", "−8.075%"],
        ["优化物理 48 层 ★", "1.532K", "+0.250K", "23.754%", "−7.502%"],
        ["旧 48 点", "1.561K", "+0.243K", "25.260%", "−9.867%"],
    ], col_widths=[Inches(2.3), Inches(1.2), Inches(1.2), Inches(1.3),
                   Inches(1.2)], font_size=10.5, header_size=10.5)
    add_textbox(s, Inches(0.58), Inches(5.7), Inches(7.2), Inches(1.2),
                "相对原 93 层：T RMSE −11.4%、|Bias| −75%；湿度持平。"
                "新基线物理含义更清晰。瓶颈：5-8km 湿度 RMSE>32%（高空信息量不足）。",
                size=11, color=C_PRIMARY_DARK, italic=True)
    # 右图
    add_image_fit(s, RESULTS / "chengdu_era5_layer48_optimized_evaluation"
                  /"04_sounding_profiles.png", Inches(8.0), Inches(1.5),
                  Inches(4.85), Inches(2.7))
    add_caption(s, Inches(8.0), Inches(4.2), Inches(4.85),
                "图10  探空独立验证廓线对比")
    add_image_fit(s, RESULTS / "chengdu_era5_layer48_optimized_evaluation"
                  /"03_sounding_height_errors.png", Inches(8.0), Inches(4.55),
                  Inches(4.85), Inches(2.4))
    add_caption(s, Inches(8.0), Inches(6.95), Inches(4.85),
                "图11  分高度误差（探空验证）")

    # ---------------- 12. 进度 / 风险 / 计划 ----------------
    s = prs.slides.add_slide(blank)
    add_header(s, "整体进度、风险与下一步计划", 12, TOTAL)
    # 进度条
    add_textbox(s, Inches(0.58), Inches(1.35), Inches(12.2), Inches(0.35),
                "▍ 当前进度总览", size=14, bold=True, color=C_PRIMARY_DARK)
    prog = [
        ("BRNN 统计反演", 100, "v4 当前最佳"),
        ("OEM 物理反演", 80, "ARTS 主后端已切换，待跑 ARTS baseline"),
        ("MonoRTM 编译", 100, "macOS + Linux 双平台"),
        ("TAPE3 光谱数据", 100, "已下载转二进制"),
        ("Sₐ 协方差", 80, "v4-derived Sₐ 已生成"),
        ("EOF/PCA 降维", 60, "基础就绪，10d 待优化"),
        ("LWC 云天 OEM", 40, "21d synthetic 实验完成"),
        ("ERA5 气压层", 5, "47/2556 天（CDS 3 天窗口）"),
        ("BRNN+OEM 桥接", 5, "待 MP-3000A Obs_BT 数据"),
    ]
    y = 1.78
    bar_w_total = Inches(7.2)
    for name, pct, note in prog:
        add_textbox(s, Inches(0.58), Inches(y), Inches(2.4), Inches(0.27),
                    name, size=10.5, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(3.0), Inches(y + 0.07), bar_w_total, Inches(0.14),
                 C_LIGHT)
        fill_color = C_ACCENT if pct >= 80 else (C_PRIMARY if pct >= 40
                                                 else C_MUTED)
        add_rect(s, Inches(3.0), Inches(y + 0.07),
                 Emu(int(bar_w_total * pct / 100)), Inches(0.14), fill_color)
        add_textbox(s, Inches(10.3), Inches(y), Inches(0.7), Inches(0.27),
                    f"{pct}%", size=10, bold=True, color=C_TEXT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(11.0), Inches(y), Inches(1.85), Inches(0.27),
                    note, size=8.5, color=C_MUTED, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.42
    # 风险 + 计划
    card(s, Inches(0.58), Inches(5.75), Inches(6.1), Inches(1.5),
         "▍ 主要风险与瓶颈", [
             "BRNN+OEM 桥接被通道不一致 + 缺 Obs_BT 阻塞。",
             "MonoRTM 单廓线 ~7s，大样本需 surrogate。",
             "高空 DOFS 低、小样本、探空时间差 1-2h。",
         ], title_fill=RGBColor(0xB0,0x3A,0x2E), b_size=10.5)
    card(s, Inches(6.83), Inches(5.75), Inches(5.95), Inches(1.5),
         "▍ 下一步（P0/P1）", [
             "P0：成都 21ch ARTS OEM baseline（n=100/200/500）。",
             "P1：Sₐ 三层递进 + BRNN 先验桥接 + EOF 10d/14d 对照。",
             "P2：NN surrogate + LWC 云天 + 不确定度校准。",
         ], title_fill=C_PRIMARY, b_size=10.5)

    # ---------------- 13. 总结 ----------------
    s = prs.slides.add_slide(blank)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(1.4), C_PRIMARY_DARK)
    add_rect(s, Emu(0), Inches(1.4), SLIDE_W, Inches(0.06), C_ACCENT)
    add_textbox(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.6),
                "总结", size=30, bold=True, color=C_WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    # 三个要点
    pts = [
        ("1", "BRNN 统计反演已达论文基准之上",
         "v4 实现 T=1.26K、RH=7.76%，六轮迭代明确了 Obs_BT 直接训练、廓线分组防泄露、Winsorize 订正等关键决策。"),
        ("2", "OEM 物理反演框架闭环完成",
         "ARTS 切换为主前向后端；MonoRTM n=100 self-consistent 基线 T=2.02K、DOFS=2.21、收敛 99%，并完成 AK/后验诊断。"),
        ("3", "成都真实数据链路打通",
         "21 通道实测亮温 + ERA5 精确配对首次实现 T/RH 同优于气候态；物理 48 层经温江探空独立验证确立为新基线。"),
    ]
    y = 1.75
    for num, title, desc in pts:
        add_rect(s, Inches(0.6), Inches(y), Inches(0.7), Inches(0.7), C_ACCENT)
        add_textbox(s, Inches(0.6), Inches(y), Inches(0.7), Inches(0.7), num,
                    size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(1.5), Inches(y - 0.05), Inches(11.2),
                    Inches(0.45), title, size=16, bold=True, color=C_PRIMARY_DARK)
        add_textbox(s, Inches(1.5), Inches(y + 0.4), Inches(11.2),
                    Inches(0.7), desc, size=12, color=C_TEXT, line_spacing=1.2)
        y += 1.35
    # 下一步一句话
    add_rect(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.95), C_LIGHT)
    add_textbox(s, Inches(0.85), Inches(6.05), Inches(11.6), Inches(0.95),
                "下一阶段重心：从「能不能做 OEM」转向提升真实物理可信度与工程可扩展性 —— "
                "ARTS baseline、Sₐ/状态向量优化、NN surrogate 加速与不确定度校准、云天 OEM。",
                size=12.5, bold=True, color=C_PRIMARY_DARK,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    add_textbox(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
                "感谢聆听  ·  欢迎指正", size=12, color=C_MUTED,
                align=PP_ALIGN.CENTER)

    prs.save(str(OUT_PATH))
    print(f"[OK] 已生成: {OUT_PATH}")
    print(f"     共 {len(prs.slides)} 页")


if __name__ == "__main__":
    build()
