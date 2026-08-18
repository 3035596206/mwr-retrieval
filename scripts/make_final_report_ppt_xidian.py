# -*- coding: utf-8 -*-
"""期末汇报 PPT 生成脚本 —— 西电空间院模板版。

以「西电空间院模板-1(1).pptx」为基底，保留其全部装饰元素（校徽、横幅、
红色顶栏、右上角标题标签、学院名）与版式，清除示例文字后，在每页的安全
内容区填入「地基微波辐射计大气温湿廓线反演」项目内容，并嵌入真实结果图。

用法:
    python scripts/make_final_report_ppt_xidian.py
产物:
    reports/汇报PPT/期末汇报_MWR温湿廓线反演_西电模板版_2026-07-29.pptx
"""

from __future__ import annotations

import copy
import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ----------------------------------------------------------------------
# 路径与常量
# ----------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parent.parent
RESULTS = PROJECT_ROOT / "results"
OUT_DIR = PROJECT_ROOT / "reports" / "汇报PPT"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "期末汇报_MWR温湿廓线反演_西电模板版_2026-07-29.pptx"

TEMPLATE = Path(r"C:\Users\Administrator\Desktop\西电空间院模板-1(1).pptx")

# 西电模板主色：暗红 C00000（来自模板）
C_RED = RGBColor(0xC0, 0x00, 0x00)
C_RED_DARK = RGBColor(0x8C, 0x00, 0x00)
C_ACCENT = RGBColor(0xE6, 0x8A, 0x00)        # 橙色强调
C_TEXT = RGBColor(0x26, 0x26, 0x26)
C_MUTED = RGBColor(0x59, 0x59, 0x59)
C_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)         # 极浅灰底
C_LIGHTRED = RGBColor(0xFB, 0xEE, 0xEE)      # 极浅红底
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE = RGBColor(0xD0, 0xCE, 0xCE)
C_HEADER = RGBColor(0x7F, 0x00, 0x00)        # 表头深红

FONT_CN = "微软雅黑"
FONT_EN = "Calibri"

# 安全区：顶部装饰占 0-1.1 英寸，右下角有装饰
CONTENT_TOP = Inches(1.25)
CONTENT_BOTTOM = Inches(7.15)
PAGE_LEFT = Inches(0.55)
PAGE_RIGHT = Inches(12.8)
CONTENT_W = PAGE_RIGHT - PAGE_LEFT


# ----------------------------------------------------------------------
# 底层辅助
# ----------------------------------------------------------------------
def _set_run(run, text, size, bold=False, color=C_TEXT, font=FONT_CN,
             italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_CN)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=C_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font=FONT_CN, line_spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        _set_run(run, line, size, bold=bold, color=color, font=font,
                 italic=italic)
    return tb


def add_bullets(slide, left, top, width, height, items, size=14,
                color=C_TEXT, line_spacing=1.25, marker_color=C_RED):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        indent = "    " * level
        marker = "● " if level == 0 else "– "
        r1 = p.add_run()
        _set_run(r1, indent + marker, size, bold=False, color=marker_color)
        r2 = p.add_run()
        _set_run(r2, text, size, bold=False, color=color)
    return tb


def add_rect(slide, left, top, width, height, fill, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def add_table(slide, left, top, width, height, data, col_widths=None,
              header_fill=C_HEADER, header_color=C_WHITE, font_size=12,
              header_size=12, zebra=True):
    rows, cols = len(data), len(data[0])
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gtable.table
    tbl = table._tbl
    for child in tbl.findall(qn("a:tblPr")):
        child.set("firstRow", "0")
        child.set("bandRow", "0")
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(20000)
            cell.margin_bottom = Emu(20000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            text = str(data[r][c])
            if r == 0:
                _set_run(run, text, header_size, bold=True, color=header_color)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                bold = (c == 0)
                color = C_TEXT
                if any(k in text for k in ("★", "当前最佳", "推荐")):
                    bold = True
                    color = C_RED
                _set_run(run, text, font_size, bold=bold, color=color)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (C_LIGHTRED if zebra and r % 2 == 1
                                            else C_WHITE)
    return gtable


def add_image_fit(slide, path, left, top, max_w, max_h, align="center",
                  valign="middle"):
    path = Path(path)
    if not path.exists():
        add_textbox(slide, left, top, max_w, max_h,
                    f"[图片缺失] {path.name}", size=12, color=C_RED,
                    anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * scale))
    h = Emu(int(ih * scale))
    x = left + (max_w - w) // 2 if align == "center" else (
        left + (max_w - w) if align == "right" else left)
    y = top + (max_h - h) // 2 if valign == "middle" else (
        top + (max_h - h) if valign == "bottom" else top)
    slide.shapes.add_picture(str(path), x, y, w, h)


def add_caption(slide, left, top, width, text):
    add_textbox(slide, left, top, width, Inches(0.28), text, size=10,
                color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


def section_title(slide, text, top=Inches(1.25)):
    """内容页的小节标题（红色竖条 + 标题文字）。"""
    add_rect(slide, PAGE_LEFT, top + Inches(0.06), Inches(0.1),
             Inches(0.42), C_RED)
    add_textbox(slide, PAGE_LEFT + Inches(0.22), top, CONTENT_W - Inches(0.3),
                Inches(0.55), text, size=20, bold=True, color=C_RED_DARK,
                anchor=MSO_ANCHOR.MIDDLE)


def card(slide, left, top, width, height, title, lines, title_color=C_WHITE,
         title_fill=C_RED, body_fill=C_LIGHTRED, t_size=13, b_size=11.5):
    bar_h = Inches(0.38)
    add_rect(slide, left, top, width, bar_h, title_fill)
    add_rect(slide, left, top + bar_h, width, height - bar_h, body_fill)
    add_textbox(slide, left + Inches(0.12), top, width - Inches(0.2), bar_h,
                title, size=t_size, bold=True, color=title_color,
                anchor=MSO_ANCHOR.MIDDLE)
    if lines:
        add_bullets(slide, left + Inches(0.14), top + bar_h + Inches(0.08),
                    width - Inches(0.24), height - bar_h - Inches(0.16),
                    lines, size=b_size, line_spacing=1.18)


# ======================================================================
# 模板操作：清除内容文字 + 设置右上角标签
# ======================================================================
def _shape_is_decoration(shape):
    """判断形状是否为模板固定装饰（应保留）：图片、或顶部横条区域内的形状。"""
    # 图片一律保留
    if shape.shape_type == 13:
        return True
    # 顶部 1.1 英寸以内的非空填充形状（横栏、装饰条）保留
    try:
        if shape.top is not None and shape.top < Inches(1.1):
            return True
    except Exception:
        pass
    # 右下角装饰
    try:
        if (shape.left is not None and shape.left > Inches(12.5)
                and shape.top is not None and shape.top > Inches(6.8)):
            return True
    except Exception:
        pass
    # 学院名（带"空间科学"字样）保留
    if shape.has_text_frame:
        if "空间科学" in shape.text_frame.text or "School of" in shape.text_frame.text:
            return True
    return False


def clear_content_shapes(slide):
    """删除非装饰形状（示例文字/目录条等），保留装饰。"""
    to_remove = []
    for shape in slide.shapes:
        if not _shape_is_decoration(shape):
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def set_corner_label(slide, text):
    """设置右上角的标签文字（'汇报'/'汇报目录'/标题等）。

    模板右上角有一个文字框承载 '汇报'/'研究内容' 等标签。找到它并改写。
    若找不到则在右上角新建。
    """
    # 寻找位于右上角区域(y<0.3, x>8.5)且非装饰文字框
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if (shape.top is not None and shape.top < Inches(0.4)
                    and shape.left is not None and shape.left > Inches(8.4)
                    and shape.left < Inches(12.5)):
                # 清空并改写
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                _set_run(run, text, 16, bold=True, color=C_WHITE)
                return shape
        except Exception:
            continue
    # 未找到则新建
    add_textbox(slide, Inches(8.9), Inches(0.18), Inches(3.8), Inches(0.5),
                text, size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    return None


def add_page_no(slide, page_no, total):
    add_textbox(slide, Inches(0.55), Inches(7.18), Inches(8), Inches(0.28),
                f"地基微波辐射计大气温湿廓线反演  ·  期末汇报", size=9,
                color=C_MUTED)
    add_textbox(slide, Inches(11.8), Inches(7.18), Inches(1.2), Inches(0.28),
                f"{page_no} / {total}", size=9, color=C_MUTED,
                align=PP_ALIGN.RIGHT)


# ======================================================================
# 主构建
# ======================================================================
def build():
    # 复制模板作为起点
    shutil.copyfile(str(TEMPLATE), str(OUT_PATH))
    prs = Presentation(str(OUT_PATH))

    # 删除模板自带的 8 页示例：需同时移除 sldIdLst 引用、presentation→slide
    # 关系、以及 slide part 本身，否则 zip 包内会残留 slide1-8.xml 与新页重名。
    sldIdLst = prs.slides._sldIdLst
    pres_part = prs.part
    rids_to_drop = []
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        if rId:
            rids_to_drop.append(rId)
        sldIdLst.remove(sldId)
    for rId in rids_to_drop:
        pres_part.drop_rel(rId)

    # 用「标题幻灯片」版式（layout[0]，模板内容页都用它）新建每一页，
    # 它继承了模板的母版背景，但没有具体装饰形状——需要手动加装饰。
    # 更稳妥：从模板原 P5（内容页）整页复制装饰。这里采用复制装饰形状法。
    # 先打开模板取一页的装饰形状做模板。
    tpl = Presentation(str(TEMPLATE))
    tpl_content_slide = tpl.slides[4]   # P5 内容页（装饰最干净）
    tpl_cover_slide = tpl.slides[0]     # P1 封面页
    tpl_toc_slide = tpl.slides[1]       # P2 目录页
    tpl_end_slide = tpl.slides[7]       # P8 结束页

    layout = prs.slide_layouts[0]       # '标题幻灯片'

    TOTAL = 12

    def new_slide_with_deco(corner_label):
        """新建一页，并复制内容页的装饰形状过来。"""
        slide = prs.slides.add_slide(layout)
        _copy_decorations(slide, tpl_content_slide)
        set_corner_label(slide, corner_label)
        return slide

    # ============================================================
    # P1 封面
    # ============================================================
    slide = prs.slides.add_slide(layout)
    _copy_decorations(slide, tpl_cover_slide)
    set_corner_label(slide, "汇报")
    # 清掉封面原有的"主题"/报告人示例文字，重新填
    _remove_text_containing(slide, ["主题"])
    _remove_text_containing(slide, ["报 告 人"])
    add_textbox(slide, Inches(0.5), Inches(2.35), Inches(8.5), Inches(0.5),
                "期末汇报 · 课题组例会", size=16,
                color=RGBColor(0x80, 0x80, 0x80))
    add_textbox(slide, Inches(0.5), Inches(2.85), Inches(8.6), Inches(1.7),
                "地基微波辐射计\n大气温湿廓线反演", size=40, bold=True,
                color=C_RED_DARK, line_spacing=1.1)
    add_textbox(slide, Inches(0.52), Inches(4.85), Inches(8.0), Inches(0.5),
                "BRNN 统计反演  ·  OEM 物理反演  ·  成都真实数据验证",
                size=15, color=C_TEXT)
    # 报告人信息块（右下）
    add_textbox(slide, Inches(8.3), Inches(4.6), Inches(4.1), Inches(2.0),
                "报 告 人：\n导    师：\n年    级：\n汇报日期：2026-07-29",
                size=14, color=C_TEXT, line_spacing=1.6)
    # 关键指标
    metrics = [("BRNN v4", "T 1.26K / RH 7.76%"),
               ("OEM MonoRTM n=100", "T 2.02K · DOFS 2.21"),
               ("成都 21ch", "真实亮温链路打通")]
    for i, (k, v) in enumerate(metrics):
        x = Inches(0.5 + i * 2.6)
        add_rect(slide, x, Inches(6.2), Inches(2.4), Inches(0.55), C_LIGHTRED)
        add_textbox(slide, x, Inches(6.22), Inches(2.4), Inches(0.27), k,
                    size=11, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(6.48), Inches(2.4), Inches(0.27), v,
                    size=10.5, color=C_TEXT, align=PP_ALIGN.CENTER)

    # ============================================================
    # P2 目录
    # ============================================================
    slide = prs.slides.add_slide(layout)
    _copy_decorations(slide, tpl_toc_slide)
    set_corner_label(slide, "汇 报 目 录")
    # 清掉模板目录里的示例项（保留"内容 CONTENTS"大字装饰）
    for kw in ["一、", "二、", "三、", "四、", "五、"]:
        _remove_text_containing(slide, [kw])
    toc = [
        ("01", "研究背景与目标"),
        ("02", "BRNN 统计反演"),
        ("03", "OEM 物理反演"),
        ("04", "成都真实数据验证"),
        ("05", "整体进度与计划"),
    ]
    for i, (no, title) in enumerate(toc):
        y = Inches(1.55 + i * 1.05)
        # 序号块
        add_rect(slide, Inches(3.7), y, Inches(0.7), Inches(0.7), C_RED)
        add_textbox(slide, Inches(3.7), y, Inches(0.7), Inches(0.7), no,
                    size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.6), y, Inches(7.0), Inches(0.7), title,
                    size=18, bold=True, color=C_TEXT,
                    anchor=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # P3 研究背景与目标
    # ============================================================
    slide = new_slide_with_deco("研究背景")
    section_title(slide, "研究背景与目标")
    card(slide, PAGE_LEFT, Inches(1.95), Inches(5.95), Inches(2.5),
         "▍ 研究背景", [
             "大气温湿廓线是数值预报、强对流与航空气象的关键输入。",
             "探空精度高但时空稀疏（每日 2 次）。",
             "微波辐射计(MWR)可实现高频连续探测。",
             "反演难点：欠定问题，垂直信息有限，云天干扰大。",
         ])
    card(slide, Inches(6.78), Inches(1.95), Inches(6.0), Inches(2.5),
         "▍ 项目目标", [
             "建立地基多通道 MWR 温度 T(z) 与湿度 RH(z) 反演系统。",
             "复现并超越论文基准（T<1.5K, RH<13%）。",
             "构建 BRNN 统计 + OEM 物理反演混合框架。",
             "在成都本地真实数据上完成验证。",
         ])
    add_table(slide, PAGE_LEFT, Inches(4.7), CONTENT_W, Inches(2.2), [
        ["数据 / 仪器", "通道配置", "覆盖", "用途"],
        ["MP-3000A（北京南郊）", "22 通道 K/V 波段", "2013–2019", "BRNN v4 训练（历史最佳）"],
        ["RPG HATPRO 配置", "14 通道（K7+V7）", "2013-01 ERA5", "OEM/MonoRTM 物理反演主线"],
        ["成都实测辐射计", "21 通道（K7+V7+W1+G5+高频1）", "2026-05", "真实数据验证"],
        ["温江探空", "483 个文件", "2026-05", "独立外部验证真值"],
    ], col_widths=[Inches(3.0), Inches(4.3), Inches(1.9), Inches(3.05)],
              font_size=11.5)

    # ============================================================
    # P4 技术路线
    # ============================================================
    slide = new_slide_with_deco("技术路线")
    section_title(slide, "技术路线：统计 + 物理混合反演")
    y_flow = Inches(1.95)
    bw, bh = Inches(2.1), Inches(1.15)
    add_rect(slide, PAGE_LEFT, y_flow, bw, bh, C_RED_DARK)
    add_textbox(slide, PAGE_LEFT, y_flow, bw, bh, "观测亮温\nObs_BT",
                size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(2.78), y_flow, bw, bh, C_RED_DARK)
    add_textbox(slide, Inches(2.78), y_flow, bw, bh, "QC / 偏差订正\n+ ERA5 标签",
                size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(5.0), Inches(1.6), Inches(2.2), Inches(0.95), C_RED)
    add_textbox(slide, Inches(5.0), Inches(1.6), Inches(2.2), Inches(0.95),
                "BRNN 统计反演\n~0.5 ms/廓线", size=13, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(5.0), Inches(2.75), Inches(2.2), Inches(0.95), C_RED)
    add_textbox(slide, Inches(5.0), Inches(2.75), Inches(2.2), Inches(0.95),
                "OEM 物理反演\n(1D-Var)", size=13, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(7.4), Inches(1.6), Inches(5.4), Inches(0.95),
             C_LIGHTRED)
    add_textbox(slide, Inches(7.55), Inches(1.6), Inches(5.1), Inches(0.95),
                "T(z), RH(z) 廓线 · 6 子模型 · 当前最佳 v4",
                size=12, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(7.4), Inches(2.75), Inches(5.4), Inches(0.95),
             C_LIGHTRED)
    add_textbox(slide, Inches(7.55), Inches(2.75), Inches(5.1), Inches(0.95),
                "T/RH + 后验误差 + AK + DOFS · ARTS 主前向模型",
                size=12, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    # 三层能力
    card(slide, PAGE_LEFT, Inches(4.05), Inches(3.95), Inches(2.85),
         "第一层 · 统计反演", [
             "BRNN v4 当前最佳",
             "T = 1.26 K，RH = 7.76%",
             "推理 ~0.5 ms/廓线",
         ], title_fill=C_RED)
    card(slide, Inches(4.69), Inches(4.05), Inches(3.95), Inches(2.85),
         "第二层 · 物理反演", [
             "OEM + ARTS/MonoRTM",
             "n=100 self-consistent",
             "DOFS ≈ 2.21，收敛 99%",
         ], title_fill=C_ACCENT)
    card(slide, Inches(8.8), Inches(4.05), Inches(3.95), Inches(2.85),
         "第三层 · 混合路线", [
             "NN 作先验 / S_a",
             "NN surrogate 加速 H(x)",
             "不确定度校准 + 云天 OEM",
         ], title_fill=C_RED_DARK)

    # ============================================================
    # P5 BRNN 方法
    # ============================================================
    slide = new_slide_with_deco("BRNN 方法")
    section_title(slide, "BRNN 统计反演：方法与 v4 关键改进")
    add_textbox(slide, PAGE_LEFT, Inches(1.95), CONTENT_W, Inches(0.35),
                "▍ 网络结构：6 个独立 BRNN 子模型（T/RH × 0-2km / 2-8km / 8-10km）",
                size=13.5, bold=True, color=C_RED_DARK)
    add_bullets(slide, Inches(0.68), Inches(2.4), CONTENT_W, Inches(1.25), [
        "输入：多通道亮温 + 地面 T2m/RH2m/Ps/IR/CLWC/CIC 等辅助特征；输出 0-10km 共 93 层 T/RH。",
        "分高度段建模：各高度区间物理特性差异大，独立网络更稳定；边界层(2km/8km)由相邻模型各预测一次。",
        "Hidden=256, Dropout=0.3, LR=1e-3, Batch=128, Early-stopping patience=20。",
    ], size=12)
    add_textbox(slide, PAGE_LEFT, Inches(3.7), CONTENT_W, Inches(0.35),
                "▍ v4 相对前序版本的关键改进（成就当前最佳）", size=13.5,
                bold=True, color=C_RED_DARK)
    items = [
        ("CLWC 廓线筛除 + K 波段逐通道 2.5σ 筛除", "剔除云天/异常亮温样本"),
        ("Winsorize 回归 BT 订正（非 OLS）", "K 波段 OMB 厚尾 → OMB std −49%"),
        ("IR_Temperature 入模", "引入红外测温作为辅助特征"),
        ("按 Profile_Index 分组划分（非时间划分）", "同廓线重复观测~5.6 次，杜绝数据泄露"),
    ]
    y0 = Inches(4.15)
    for i, (t, d) in enumerate(items):
        col = i % 2
        row = i // 2
        x = PAGE_LEFT + Inches(col * 6.1)
        yy = y0 + Inches(row * 1.05)
        add_rect(slide, x, yy, Inches(0.12), Inches(0.9), C_RED)
        add_textbox(slide, x + Inches(0.25), yy, Inches(5.7), Inches(0.45), t,
                    size=12.5, bold=True, color=C_TEXT)
        add_textbox(slide, x + Inches(0.25), yy + Inches(0.46), Inches(5.7),
                    Inches(0.4), "→ " + d, size=11, color=C_MUTED)
    add_textbox(slide, PAGE_LEFT, Inches(6.4), CONTENT_W, Inches(0.5),
                "核心结论：v3 证实 Sim→Obs domain gap 退化 3.7×，故 v4 改为 Obs_BT 直接训练，"
                "成为当前最优统计反演模型。",
                size=12, color=C_RED_DARK, italic=True)

    # ============================================================
    # P6 BRNN 结果
    # ============================================================
    slide = new_slide_with_deco("BRNN 结果")
    section_title(slide, "BRNN v4 反演结果（MP-3000A）")
    add_textbox(slide, PAGE_LEFT, Inches(1.95), Inches(4.3), Inches(0.35),
                "▍ 核心指标", size=13, bold=True, color=C_RED_DARK)
    add_table(slide, PAGE_LEFT, Inches(2.35), Inches(4.35), Inches(1.65), [
        ["指标", "v4", "论文基准"],
        ["T RMSE", "1.26 K", "<1.5 K"],
        ["RH RMSE", "7.76%", "<13%"],
        ["推理速度", "~0.5 ms/廓线", "—"],
    ], col_widths=[Inches(1.5), Inches(1.45), Inches(1.4)], font_size=12)
    add_textbox(slide, PAGE_LEFT, Inches(4.25), Inches(4.3), Inches(0.35),
                "▍ 分层精度", size=13, bold=True, color=C_RED_DARK)
    add_table(slide, PAGE_LEFT, Inches(4.65), Inches(4.35), Inches(2.05), [
        ["高度区间", "T RMSE", "RH RMSE"],
        ["0–0.5 km", "1.09 K", "5.78%"],
        ["2–8 km", "1.30 K", "9.50%"],
        ["5–10 km", "1.50 K", "10.31%"],
    ], col_widths=[Inches(1.7), Inches(1.35), Inches(1.3)], font_size=11.5)
    # 右图
    add_image_fit(slide, RESULTS / "T_scatter.png", Inches(5.2), Inches(1.95),
                  Inches(3.75), Inches(2.5))
    add_caption(slide, Inches(5.2), Inches(4.45), Inches(3.75),
                "图1  温度散点（v4）")
    add_image_fit(slide, RESULTS / "RH_scatter.png", Inches(9.1),
                  Inches(1.95), Inches(3.75), Inches(2.5))
    add_caption(slide, Inches(9.1), Inches(4.45), Inches(3.75),
                "图2  相对湿度散点（v4）")
    add_image_fit(slide, RESULTS / "T_error_profile.png", Inches(5.2),
                  Inches(4.85), Inches(3.75), Inches(2.2))
    add_caption(slide, Inches(5.2), Inches(7.0), Inches(3.75),
                "图3  温度误差廓线")
    add_image_fit(slide, RESULTS / "RH_error_profile.png", Inches(9.1),
                  Inches(4.85), Inches(3.75), Inches(2.2))
    add_caption(slide, Inches(9.1), Inches(7.0), Inches(3.75),
                "图4  湿度误差廓线")

    # ============================================================
    # P7 BRNN 迭代历程
    # ============================================================
    slide = new_slide_with_deco("迭代历程")
    section_title(slide, "BRNN 迭代历程与关键决策")
    add_textbox(slide, PAGE_LEFT, Inches(1.95), CONTENT_W, Inches(0.35),
                "▍ 六轮迭代：从数据泄露基线到当前最佳 v4", size=13.5,
                bold=True, color=C_RED_DARK)
    add_table(slide, PAGE_LEFT, Inches(2.35), CONTENT_W, Inches(2.95), [
        ["版本", "方案", "T RMSE", "RH RMSE", "结论"],
        ["v1", "原始 Obs_BT，时间划分", "3.03 K", "16.55%", "基线，存在数据泄露与质量问题"],
        ["v2", "Obs_BT 订正 + 廓线分组", "1.45 K", "9.03%", "首个强结果"],
        ["v3", "Sim_BT 训练，Obs_BT 测试", "2.65 K", "12.00%", "暴露 Sim→Obs domain gap"],
        ["v4 ★", "Obs_BT 订正 + K波段过滤 + IR入模", "1.26 K", "7.76%", "当前最佳"],
        ["v6", "Sim→Obs 两阶段训练", "1.92 K", "10.34%", "Sim 路线有进展但未超 v4"],
    ], col_widths=[Inches(1.1), Inches(4.6), Inches(1.5), Inches(1.5),
                   Inches(3.55)], font_size=11.5)
    add_textbox(slide, PAGE_LEFT, Inches(5.5), CONTENT_W, Inches(0.35),
                "▍ 关键决策记录", size=13.5, bold=True, color=C_RED_DARK)
    decisions = [
        ("Obs_BT 直接训练（v4）而非 Sim_BT", "v3 证明 Sim→Obs 退化 3.7×"),
        ("6 个独立 BRNN 而非单一多输出", "各高度段物理特性差异大"),
        ("廓线分组划分而非时间划分", "同廓线重复观测~5.6次，必须防泄露"),
        ("Winsorize 回归订正而非 OLS", "K 波段 OMB 厚尾分布"),
    ]
    for i, (t, d) in enumerate(decisions):
        col = i % 2
        row = i // 2
        x = PAGE_LEFT + Inches(col * 6.1)
        yy = Inches(5.95 + row * 0.72)
        add_textbox(slide, x, yy, Inches(5.9), Inches(0.35),
                    "▸ " + t, size=12, bold=True, color=C_TEXT)
        add_textbox(slide, x + Inches(0.3), yy + Inches(0.32), Inches(5.6),
                    Inches(0.35), d, size=10.5, color=C_MUTED)

    # ============================================================
    # P8 OEM 框架
    # ============================================================
    slide = new_slide_with_deco("OEM 框架")
    section_title(slide, "OEM 物理反演框架")
    add_textbox(slide, PAGE_LEFT, Inches(1.95), CONTENT_W, Inches(0.35),
                "▍ 最优估计 (OEM / 1D-Var) 目标函数", size=13.5, bold=True,
                color=C_RED_DARK)
    add_rect(slide, PAGE_LEFT, Inches(2.4), CONTENT_W, Inches(0.7), C_LIGHTRED)
    add_textbox(slide, PAGE_LEFT, Inches(2.4), CONTENT_W, Inches(0.7),
                "J(x) = (x−xₐ)ᵀ Sₐ⁻¹ (x−xₐ)  +  (y−H(x))ᵀ Sₑ⁻¹ (y−H(x))",
                size=16, bold=True, color=C_RED_DARK, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    add_table(slide, PAGE_LEFT, Inches(3.3), Inches(7.4), Inches(3.3), [
        ["符号", "含义", "当前实现"],
        ["x", "状态向量", "14d T7+RH7（→ 21d 含 LWC）"],
        ["xₐ", "背景场", "ERA5 扰动（→ BRNN first guess）"],
        ["Sₐ", "背景误差协方差", "指数相关 / v4-derived"],
        ["y", "观测亮温", "synthetic / MonoRTM / Obs_BT"],
        ["H(x)", "前向模型", "ARTS 主线；MonoRTM 历史"],
        ["Sₑ", "观测误差协方差", "K-band 1.5K, V-band 0.5K"],
    ], col_widths=[Inches(1.1), Inches(2.4), Inches(3.9)], font_size=11)
    card(slide, Inches(8.18), Inches(3.3), Inches(4.6), Inches(3.3),
         "▍ 已实现的 OEM 能力", [
             "LM / Gauss-Newton 求解器",
             "有限差分 Jacobian",
             "状态向量打包/解包",
             "Averaging kernel / DOFS",
             "后验协方差诊断",
             "self-consistent 闭环验证",
         ], title_fill=C_ACCENT, b_size=11.5)
    add_textbox(slide, PAGE_LEFT, Inches(6.75), CONTENT_W, Inches(0.4),
                "状态向量粗分层：0-0.5 / 0.5-1 / 1-2 / 2-3 / 3-5 / 5-8 / 8-10 km，"
                "T7+RH7 共 14 维；后续扩展 10d EOF 与 21d 云天状态。",
                size=10.5, color=C_MUTED, italic=True)

    # ============================================================
    # P9 OEM 结果
    # ============================================================
    slide = new_slide_with_deco("OEM 结果")
    section_title(slide, "OEM 实验结果：MonoRTM n=100 基线")
    add_textbox(slide, PAGE_LEFT, Inches(1.95), Inches(7.2), Inches(0.35),
                "▍ 2013-01 三类 POC + n=100 基线", size=13.5, bold=True,
                color=C_RED_DARK)
    add_table(slide, PAGE_LEFT, Inches(2.35), Inches(7.2), Inches(2.25), [
        ["实验（前向/样本）", "T RMSE", "RH RMSE", "BT RMS", "DOFS"],
        ["Self-consistent simple n=20", "2.64→1.95K", "6.43→6.06%", "1.59→0.54K", "2.10"],
        ["Forward-mismatch simple n=20", "→27.0K", "7.59→7.22%", "→20.2K", "2.21"],
        ["MonoRTM self-consistent n=100", "2.64→2.02K", "6.49→6.20%", "5.03→0.61K", "2.21"],
    ], col_widths=[Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.1),
                   Inches(0.7)], font_size=10.5, header_size=10.5)
    add_textbox(slide, PAGE_LEFT, Inches(4.75), Inches(7.2), Inches(0.35),
                "▍ 核心结论", size=12.5, bold=True, color=C_RED_DARK)
    add_bullets(slide, Inches(0.68), Inches(5.15), Inches(7.0), Inches(2.0), [
        "self-consistent 验证算法链路可用，收敛率 99%、平均迭代 7.79。",
        "forward-mismatch 证明 H(x) 与观测不一致会导致 T 严重退化 → 必须先闭环验证。",
        "MonoRTM 是当前最可信物理 POC；BT residual 降幅最大（5.03→0.61K）。",
        "DOFS≈2.21/14，符合地基 MWR 垂直信息量有限的预期。",
    ], size=11.5)
    add_image_fit(slide,
                  RESULTS / "oem_201301_self_consistent_monortm_n100"
                  / "rmse_profiles.png",
                  Inches(8.0), Inches(2.0), Inches(4.85), Inches(2.55))
    add_caption(slide, Inches(8.0), Inches(4.55), Inches(4.85),
                "图5  MonoRTM n=100：Prior→Posterior RMSE 廓线")
    add_image_fit(slide,
                  RESULTS / "oem_201301_self_consistent_monortm_n100"
                  / "bt_dofs.png",
                  Inches(8.0), Inches(4.9), Inches(4.85), Inches(2.1))
    add_caption(slide, Inches(8.0), Inches(7.0), Inches(4.85),
                "图6  BT 残差收敛 与 DOFS 分布")

    # ============================================================
    # P10 成都验证
    # ============================================================
    slide = new_slide_with_deco("成都验证")
    section_title(slide, "成都真实数据验证")
    add_textbox(slide, PAGE_LEFT, Inches(1.95), Inches(6.1), Inches(0.35),
                "▍ ERA5 精确配对 + 物理 48 层（温江探空独立验证）",
                size=13, bold=True, color=C_RED_DARK)
    add_table(slide, PAGE_LEFT, Inches(2.35), Inches(6.1), Inches(1.75), [
        ["数据集划分", "样本 / 天数"],
        ["实测亮温总记录", "163 条"],
        ["与 ERA5 同 UTC 配对", "139 条 / 16 天"],
        ["Train / Val / Test", "69 / 22 / 48 条"],
    ], col_widths=[Inches(3.4), Inches(2.7)], font_size=11.5)
    add_textbox(slide, PAGE_LEFT, Inches(4.25), Inches(6.1), Inches(0.35),
                "▍ 关键结果", size=12.5, bold=True, color=C_RED_DARK)
    add_table(slide, PAGE_LEFT, Inches(4.65), Inches(6.1), Inches(2.15), [
        ["方案 / 变量", "T RMSE", "RH RMSE"],
        ["ERA5 Hybrid（Ridge+BRNN）", "1.479 K", "20.966%"],
        ["物理 48 层（探空验证）", "1.532 K", "23.754%"],
        ["相对气候态基线改善", "35% / —", "13% / 持平"],
    ], col_widths=[Inches(3.1), Inches(1.5), Inches(1.5)], font_size=11)
    add_textbox(slide, PAGE_LEFT, Inches(6.9), Inches(6.1), Inches(0.4),
                "首次实现 T/RH 同时优于气候态；瓶颈：5-8km 湿度 RMSE>32%。",
                size=10.5, color=C_RED_DARK, italic=True)
    # 右图
    add_image_fit(slide, RESULTS / "chengdu_era5_figures"
                  / "03_prediction_scatter.png", Inches(6.95), Inches(2.05),
                  Inches(5.9), Inches(2.5))
    add_caption(slide, Inches(6.95), Inches(4.55), Inches(5.9),
                "图7  成都 21ch Hybrid 预测散点（T/RH）")
    add_image_fit(slide,
                  RESULTS / "chengdu_era5_layer48_optimized_evaluation"
                  / "04_sounding_profiles.png",
                  Inches(6.95), Inches(4.9), Inches(5.9), Inches(2.1))
    add_caption(slide, Inches(6.95), Inches(7.0), Inches(5.9),
                "图8  物理 48 层温江探空独立验证廓线")

    # ============================================================
    # P11 进度与计划
    # ============================================================
    slide = new_slide_with_deco("进度与计划")
    section_title(slide, "整体进度、风险与下一步计划")
    add_textbox(slide, PAGE_LEFT, Inches(1.95), CONTENT_W, Inches(0.35),
                "▍ 当前进度总览", size=13.5, bold=True, color=C_RED_DARK)
    prog = [
        ("BRNN 统计反演", 100, "v4 当前最佳"),
        ("OEM 物理反演", 80, "ARTS 主后端已切换"),
        ("MonoRTM 编译", 100, "macOS + Linux 双平台"),
        ("Sₐ 协方差", 80, "v4-derived Sₐ 已生成"),
        ("EOF/PCA 降维", 60, "基础就绪，10d 待优化"),
        ("LWC 云天 OEM", 40, "21d synthetic 完成"),
        ("ERA5 气压层", 5, "47/2556 天"),
        ("BRNN+OEM 桥接", 5, "待 Obs_BT 数据"),
    ]
    y = 2.4
    bar_total = Inches(7.0)
    for name, pct, note in prog:
        add_textbox(slide, PAGE_LEFT, Inches(y), Inches(2.3), Inches(0.26),
                    name, size=10.5, color=C_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(2.95), Inches(y + 0.07), bar_total,
                 Inches(0.13), C_LIGHT)
        fc = C_RED if pct >= 80 else (C_ACCENT if pct >= 40 else C_MUTED)
        add_rect(slide, Inches(2.95), Inches(y + 0.07),
                 Emu(int(bar_total * pct / 100)), Inches(0.13), fc)
        add_textbox(slide, Inches(10.05), Inches(y), Inches(0.7),
                    Inches(0.26), f"{pct}%", size=10, bold=True, color=C_TEXT,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(10.8), Inches(y), Inches(2.0),
                    Inches(0.26), note, size=8.5, color=C_MUTED,
                    anchor=MSO_ANCHOR.MIDDLE)
        y += 0.36
    card(slide, PAGE_LEFT, Inches(5.55), Inches(6.1), Inches(1.55),
         "▍ 主要风险与瓶颈", [
             "BRNN+OEM 桥接被通道不一致 + 缺 Obs_BT 阻塞。",
             "MonoRTM 单廓线 ~7s，大样本需 surrogate。",
             "高空 DOFS 低、小样本、探空时间差 1-2h。",
         ], title_fill=RGBColor(0xA6, 0x2A, 0x2A), b_size=10.5)
    card(slide, Inches(6.83), Inches(5.55), Inches(5.95), Inches(1.55),
         "▍ 下一步（P0/P1）", [
             "P0：成都 21ch ARTS OEM baseline（n=100/200/500）。",
             "P1：Sₐ 三层递进 + BRNN 先验桥接 + EOF 对照。",
             "P2：NN surrogate + LWC 云天 + 不确定度校准。",
         ], title_fill=C_RED, b_size=10.5)

    # ============================================================
    # P12 总结
    # ============================================================
    slide = prs.slides.add_slide(layout)
    _copy_decorations(slide, tpl_end_slide)
    set_corner_label(slide, "总结")
    # 清掉结束页示例文字
    _remove_text_containing(slide, ["感谢各位老师"])
    _remove_text_containing(slide, ["恳请指正"])
    add_textbox(slide, PAGE_LEFT, Inches(1.4), CONTENT_W, Inches(0.6),
                "总结", size=28, bold=True, color=C_RED_DARK)
    pts = [
        ("1", "BRNN 统计反演已达论文基准之上",
         "v4 实现 T=1.26K、RH=7.76%，六轮迭代明确了 Obs_BT 直接训练、廓线分组防泄露、Winsorize 订正等关键决策。"),
        ("2", "OEM 物理反演框架闭环完成",
         "ARTS 切换为主前向后端；MonoRTM n=100 self-consistent 基线 T=2.02K、DOFS=2.21、收敛 99%，并完成 AK/后验诊断。"),
        ("3", "成都真实数据链路打通",
         "21 通道实测亮温 + ERA5 精确配对首次实现 T/RH 同优于气候态；物理 48 层经温江探空独立验证确立为新基线。"),
    ]
    y = 2.25
    for num, title, desc in pts:
        add_rect(slide, PAGE_LEFT, Inches(y), Inches(0.7), Inches(0.7), C_RED)
        add_textbox(slide, PAGE_LEFT, Inches(y), Inches(0.7), Inches(0.7), num,
                    size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.5), Inches(y - 0.05), Inches(11.2),
                    Inches(0.45), title, size=16, bold=True, color=C_RED_DARK)
        add_textbox(slide, Inches(1.5), Inches(y + 0.4), Inches(11.2),
                    Inches(0.7), desc, size=12, color=C_TEXT, line_spacing=1.2)
        y += 1.35
    add_rect(slide, PAGE_LEFT, Inches(6.35), Inches(12.25), Inches(0.55),
             C_LIGHTRED)
    add_textbox(slide, Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.55),
                "下一阶段重心：ARTS baseline、Sₐ/状态向量优化、NN surrogate 加速、不确定度校准与云天 OEM。",
                size=12, bold=True, color=C_RED_DARK, anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.CENTER)

    # 给 P3-P12 加页码（P1封面、P2目录、P12总结不加）
    for idx in range(2, 11):
        add_page_no(prs.slides[idx], idx + 1, TOTAL)

    prs.save(str(OUT_PATH))
    print(f"[OK] 已生成: {OUT_PATH}")
    print(f"     共 {len(prs.slides)} 页")


# ----------------------------------------------------------------------
# 装饰形状复制 & 文字清理
# ----------------------------------------------------------------------
def _copy_decorations(target_slide, source_slide):
    """把 source_slide 的装饰形状复制到 target_slide。

    复制策略：用 lxml 深拷贝源页的 sp/pic 元素，并重建图片关系。
    为简单稳妥，这里改为：复制源页 *所有* 形状的 XML 到目标页，
    但跳过纯示例文字框（含特定关键词的）。
    """
    # 源页 XML 根的 spTree
    src_spTree = source_slide.shapes._spTree
    tgt_spTree = target_slide.shapes._spTree
    # 收集目标页已有形状（避免重复），新页一开始只有 layout 占位符
    # 复制源页的 sp 和 pic（装饰 + 内容），交给后续 clear 步骤决定删哪些
    # 但为避免与 add_slide 自带的 layout placeholders 冲突，只复制非占位符形状
    copied = 0
    for elem in list(src_spTree):
        tag = elem.tag
        # 只处理 sp / pic / grpSp / cxnSp（图形类），跳过 nvGrpSpPr/grpSpPr/custXml
        local = tag.split("}")[-1] if "}" in tag else tag
        if local in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            # 跳过占位符 sp（layout 继承来的）
            is_placeholder = False
            if local == "sp":
                nvSpPr = elem.find(qn("p:nvSpPr"))
                if nvSpPr is not None:
                    nvPr = nvSpPr.find(qn("p:nvPr"))
                    if nvPr is not None and nvPr.find(qn("p:ph")) is not None:
                        is_placeholder = True
            if is_placeholder:
                continue
            # 深拷贝
            new_elem = copy.deepcopy(elem)
            # 若是 pic，需重定向 r:embed 到目标页的 media（重新添加图片）
            if local == "pic":
                _relink_picture(target_slide, source_slide, new_elem)
            tgt_spTree.append(new_elem)
            copied += 1
    return copied


def _relink_picture(target_slide, source_slide, pic_elem):
    """把 pic 元素里的 r:embed(rId) 重指向目标页新添加的同一图片文件。"""
    # 找 blip embed
    blipFill = pic_elem.find(qn("p:blipFill"))
    if blipFill is None:
        return
    blip = blipFill.find(qn("a:blip"))
    if blip is None:
        return
    embed_attr = qn("r:embed")
    old_rid = blip.get(embed_attr)
    if not old_rid:
        return
    # 从源 slide part 拿到 image part
    try:
        src_part = source_slide.part
        image_part = src_part.related_part(old_rid)
        # 在目标 slide part 添加同一 image part，得到新 rId
        tgt_part = target_slide.part
        new_rid = tgt_part.relate_to(image_part,
                                     "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
        blip.set(embed_attr, new_rid)
    except Exception as e:
        print(f"[warn] 图片重链接失败: {e}")


def _remove_text_containing(slide, keywords):
    """删除文本包含任一关键词的形状（用于清示例文字）。"""
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if any(kw in txt for kw in keywords):
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


if __name__ == "__main__":
    build()
